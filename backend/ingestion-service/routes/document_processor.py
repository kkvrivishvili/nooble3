"""
Procesamiento de diferentes tipos de documentos.
"""

import logging
import mimetypes
import os
from typing import Dict, Any, List, Optional, BinaryIO, Union

from fastapi import UploadFile
import httpx

from common.errors import (
    ServiceError, DocumentProcessingError, 
    ValidationError, TextTooLargeError
)

from config import get_settings

logger = logging.getLogger(__name__)

# Configuración
settings = get_settings()
MAX_FILE_SIZE_MB = settings.max_file_size_mb
SUPPORTED_FILE_TYPES = settings.supported_file_types

async def validate_file(file: UploadFile, max_size_mb: int = MAX_FILE_SIZE_MB) -> Dict[str, Any]:
    """
    Valida un archivo cargado, verificando tipo y tamaño.
    
    Args:
        file: Archivo a validar
        max_size_mb: Tamaño máximo permitido en MB
        
    Returns:
        Dict: Información del archivo
        
    Raises:
        ValidationError: Si el archivo no es válido
    """
    if not file.filename:
        raise ValidationError(
            message="Nombre de archivo inválido",
            details={"filename": None}
        )
    
    # Verificar extensión
    file_extension = os.path.splitext(file.filename)[1].lower().replace(".", "")
    if file_extension not in SUPPORTED_FILE_TYPES:
        raise ValidationError(
            message=f"Tipo de archivo no soportado: {file_extension}",
            details={
                "extension": file_extension, 
                "supported_types": SUPPORTED_FILE_TYPES
            }
        )
    
    # Verificar tipo MIME
    content_type = file.content_type or mimetypes.guess_type(file.filename)[0] or "application/octet-stream"
    
    # Verificar tamaño (aproximado)
    try:
        # Verificar posición actual y tamaño
        await file.seek(0, 2)  # Ir al final
        size_bytes = await file.tell()  # Obtener posición
        await file.seek(0)  # Volver al inicio
        
        size_mb = size_bytes / (1024 * 1024)
        if size_mb > max_size_mb:
            raise ValidationError(
                message=f"Archivo demasiado grande ({size_mb:.2f} MB). Máximo: {max_size_mb} MB",
                details={
                    "file_size_mb": size_mb, 
                    "max_size_mb": max_size_mb
                }
            )
        
        return {
            "filename": file.filename,
            "size": size_bytes,
            "type": file_extension,
            "content_type": content_type
        }
    except Exception as e:
        if isinstance(e, ValidationError):
            raise
        raise ValidationError(
            message=f"Error validando archivo: {str(e)}",
            details={"filename": file.filename}
        )

async def process_text(text: str) -> str:
    """
    Procesa texto plano.
    
    Args:
        text: Texto a procesar
        
    Returns:
        str: Texto procesado
    """
    # Procesar texto (eliminar caracteres no deseados, normalizar, etc.)
    return text.strip()

async def process_pdf(file_content: bytes) -> str:
    """
    Procesa un archivo PDF y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del PDF
    """
    import PyPDF2
    import io
    
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n\n"
        
        return text
    except Exception as e:
        logger.error(f"Error procesando PDF: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando PDF: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_docx(file_content: bytes) -> str:
    """
    Procesa un archivo Word (DOCX) y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    import docx2txt
    import io
    
    try:
        docx_file = io.BytesIO(file_content)
        text = docx2txt.process(docx_file)
        return text
    except Exception as e:
        logger.error(f"Error procesando DOCX: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando DOCX: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_xlsx(file_content: bytes) -> str:
    """
    Procesa un archivo Excel (XLSX) y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    import pandas as pd
    import io
    
    try:
        xlsx_file = io.BytesIO(file_content)
        # Leer todas las hojas
        dfs = pd.read_excel(xlsx_file, sheet_name=None)
        
        text = ""
        for sheet_name, df in dfs.items():
            text += f"Hoja: {sheet_name}\n\n"
            text += df.to_string(index=False) + "\n\n"
        
        return text
    except Exception as e:
        logger.error(f"Error procesando XLSX: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando XLSX: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_pptx(file_content: bytes) -> str:
    """
    Procesa un archivo PowerPoint (PPTX) y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    from pptx import Presentation
    import io
    
    try:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text = ""
        for i, slide in enumerate(presentation.slides):
            text += f"Diapositiva {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        return text
    except Exception as e:
        logger.error(f"Error procesando PPTX: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando PPTX: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_csv(file_content: bytes) -> str:
    """
    Procesa un archivo CSV y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    import pandas as pd
    import io
    
    try:
        csv_file = io.BytesIO(file_content)
        df = pd.read_csv(csv_file)
        text = df.to_string(index=False)
        return text
    except Exception as e:
        logger.error(f"Error procesando CSV: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando CSV: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_html(file_content: bytes) -> str:
    """
    Procesa un archivo HTML y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    from bs4 import BeautifulSoup
    
    try:
        soup = BeautifulSoup(file_content, 'html.parser')
        # Eliminar scripts y estilos
        for script in soup(["script", "style"]):
            script.extract()
        
        text = soup.get_text()
        # Limpiar espacios en blanco
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except Exception as e:
        logger.error(f"Error procesando HTML: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando HTML: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_markdown(file_content: bytes) -> str:
    """
    Procesa un archivo Markdown y extrae su texto.
    
    Args:
        file_content: Contenido del archivo en bytes
        
    Returns:
        str: Texto extraído del documento
    """
    try:
        text = file_content.decode('utf-8')
        # Opcional: convertir a texto plano eliminando marcas markdown
        # import markdown
        # html = markdown.markdown(text)
        # from bs4 import BeautifulSoup
        # soup = BeautifulSoup(html, features='html.parser')
        # text = soup.get_text()
        
        return text
    except Exception as e:
        logger.error(f"Error procesando Markdown: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando Markdown: {str(e)}",
            details={"error_type": type(e).__name__}
        )

async def process_url(url: str) -> str:
    """
    Procesa una URL y extrae su contenido.
    
    Args:
        url: URL a procesar
        
    Returns:
        str: Texto extraído del contenido
    """
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url)
            response.raise_for_status()
            
            content_type = response.headers.get("content-type", "").lower()
            
            if "text/html" in content_type:
                # Procesar como HTML
                return await process_html(response.content)
            elif "application/pdf" in content_type:
                # Procesar como PDF
                return await process_pdf(response.content)
            elif "text/plain" in content_type:
                # Procesar como texto
                return response.text
            else:
                # Intentar procesar como texto por defecto
                return response.text
    except Exception as e:
        logger.error(f"Error procesando URL {url}: {str(e)}")
        raise DocumentProcessingError(
            message=f"Error procesando URL: {str(e)}",
            details={"url": url}
        )

async def process_file(
    file_content: bytes, 
    file_type: str,
    metadata: Dict[str, Any]
) -> str:
    """
    Procesa un archivo según su tipo.
    
    Args:
        file_content: Contenido del archivo en bytes
        file_type: Tipo del archivo (pdf, docx, etc.)
        metadata: Metadatos del archivo
        
    Returns:
        str: Texto extraído del archivo
    """
    file_type = file_type.lower()
    
    try:
        if file_type == "pdf":
            return await process_pdf(file_content)
        elif file_type == "docx":
            return await process_docx(file_content)
        elif file_type == "xlsx":
            return await process_xlsx(file_content)
        elif file_type == "pptx":
            return await process_pptx(file_content)
        elif file_type == "csv":
            return await process_csv(file_content)
        elif file_type == "html":
            return await process_html(file_content)
        elif file_type == "md":
            return await process_markdown(file_content)
        elif file_type == "txt":
            return await process_text(file_content.decode('utf-8'))
        else:
            raise ValidationError(
                message=f"Tipo de archivo no soportado: {file_type}",
                details={
                    "file_type": file_type,
                    "supported_types": SUPPORTED_FILE_TYPES
                }
            )
    except Exception as e:
        logger.error(f"Error procesando archivo {file_type}: {str(e)}")
        if isinstance(e, ServiceError):
            raise
        raise DocumentProcessingError(
            message=f"Error procesando archivo {file_type}: {str(e)}",
            details={
                "file_type": file_type,
                "metadata": metadata
            }
        )